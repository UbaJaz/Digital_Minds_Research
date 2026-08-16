"""Build the submission slide deck as a .pptx — six slides.

Same content and visual language as presentation.html, as a real PowerPoint file so it
can be opened, edited and presented locally.

The palette and typography are carried over from the team's previous hackathon deck
(AfriGuard, June 2026): dark ground, Calibri, one amber accent for figures, blue for the
external observer, red reserved for caveats. Report figures are light-background PNGs and
are placed on white cards deliberately — they read as printed exhibits, not as chrome.

Every number on these slides comes from 10_report.md. Nothing is computed here.

Run:  .venv/Scripts/python scripts/build_pptx.py  ->  Digital_Minds_Track3_Slides.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "figures"
OUT = ROOT / "Digital_Minds_Track3_Slides.pptx"

# ---- palette (carried from the AfriGuard deck) ----------------------------------------
GROUND = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x22)
PANEL2 = RGBColor(0x1A, 0x1F, 0x2E)
RULE = RGBColor(0x2A, 0x32, 0x3C)
INK = RGBColor(0xE9, 0xEE, 0xF4)
MUTED = RGBColor(0x98, 0xA3, 0xB0)
DIM = RGBColor(0x6B, 0x76, 0x83)
AMBER = RGBColor(0xF0, 0xC0, 0x40)
BLUE = RGBColor(0x4A, 0x90, 0xD9)
RED = RGBColor(0xE2, 0x4A, 0x4A)
GREEN = RGBColor(0x4A, 0x9B, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
BODY_W = W - MARGIN * 2


# ---- primitives -----------------------------------------------------------------------
def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs: Presentation, ground: RGBColor = GROUND):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = ground
    return s


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size, font=SANS, color=INK, bold=False, space_after=8,
         first=False, align=PP_ALIGN.LEFT, line=None, caps=False, spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if line:
        p.line_spacing = line
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text.upper() if caps else text
    f = r.font
    f.size, f.name, f.bold = Pt(size), font, bold
    f.color.rgb = color
    if spacing is not None:  # letter-spacing via raw XML (python-pptx has no API for it)
        r.font._rPr.set("spc", str(int(spacing * 100)))
    return p


def rich(tf, runs, *, first=False, space_after=8, line=1.28, align=PP_ALIGN.LEFT):
    """One paragraph, several differently-styled runs: (text, size, colour, bold, font)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    p.space_after = Pt(space_after)
    for text, size, color, bold, font in runs:
        r = p.add_run()
        r.text = text
        r.font.size, r.font.name, r.font.bold = Pt(size), font, bold
        r.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill=PANEL, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def eyebrow(slide, text, y=Inches(0.52)):
    tf = textbox(slide, MARGIN, y, BODY_W, Inches(0.34))
    para(tf, text, size=11.5, font=MONO, color=AMBER, first=True, caps=True, spacing=1.6)


def heading(slide, text, y=Inches(0.95), size=34, color=INK, w=None):
    tf = textbox(slide, MARGIN, y, w or BODY_W, Inches(1.1))
    para(tf, text, size=size, color=color, bold=True, first=True, line=1.06)
    return tf


def footer(slide, num, label):
    """Slide number and section label on a hairline — a linear talk, so ordinals inform."""
    ln = rect(slide, MARGIN, H - Inches(0.62), BODY_W, Emu(9525), fill=RULE)
    ln.shadow.inherit = False
    tf = textbox(slide, MARGIN, H - Inches(0.5), BODY_W, Inches(0.3))
    rich(tf, [(f"{num}", 10.5, AMBER, True, MONO),
              (f"  /  6     {label}", 10.5, DIM, False, MONO)], first=True, space_after=0)


def picture(slide, name, y, height, cx=None, pad=Inches(0.14)):
    """Light-background report figure, matted on a white card so it reads as an exhibit."""
    img = FIG / name
    pic = slide.shapes.add_picture(str(img), Inches(0), y, height=height)
    left = int(((cx or W) - pic.width) / 2) if cx is None else int(cx - pic.width / 2)
    pic.left = left
    card = rect(slide, left - pad, y - pad, pic.width + pad * 2, pic.height + pad * 2, fill=WHITE)
    card.shadow.inherit = False
    # add_shape draws above the picture; push the card behind it
    sp = card._element
    sp.getparent().remove(sp)
    pic._element.addprevious(sp)
    return pic


# ---- slides ---------------------------------------------------------------------------
def slide_title(prs):
    s = blank(prs)
    rect(s, Inches(0), Inches(0), W, Inches(0.09), fill=AMBER)

    tf = textbox(s, MARGIN, Inches(0.72), BODY_W, Inches(0.3))
    para(tf, "Digital Minds Research Sprint  ·  Track 3  ·  Introspection & Self-Report Reliability",
         size=12, font=MONO, color=MUTED, first=True, spacing=1.0)

    tf = textbox(s, MARGIN, Inches(1.55), Inches(11.4), Inches(2.2))
    para(tf, "Beaten by a Cheap Surface Classifier", size=50, color=INK, bold=True, first=True,
         line=1.02, space_after=4)
    para(tf, "A Capability-Controlled Test of Privileged Self-Access", size=30, color=AMBER,
         line=1.08)

    bar = rect(s, MARGIN, Inches(3.72), Inches(2.6), Emu(19050), fill=RULE)
    bar.shadow.inherit = False

    tf = textbox(s, MARGIN, Inches(4.05), Inches(9.6), Inches(1.0))
    para(tf, "Does behavioural self-prediction provide information unavailable to an "
             "external observer?", size=21, color=MUTED, first=True, line=1.28)

    # authors
    for i, (name, aff) in enumerate([
        ("Ubayd Hattas", "Computer Science, Statistics & Data Science\nUniversity of Cape Town"),
        ("Jaswin Chinthala", "Electrical Engineering\nUniversity of Cape Town"),
    ]):
        x = MARGIN + i * Inches(4.3)
        tf = textbox(s, x, Inches(5.55), Inches(4.0), Inches(1.1))
        para(tf, name, size=16, color=INK, bold=True, first=True, space_after=3)
        for j, ln in enumerate(aff.split("\n")):
            para(tf, ln, size=12, color=MUTED, space_after=1, line=1.2)

    tf = textbox(s, W - MARGIN - Inches(3.6), Inches(5.55), Inches(3.6), Inches(1.1))
    para(tf, "9,269 scored trials  ·  $3.12", size=13, font=MONO, color=AMBER, bold=True,
         first=True, align=PP_ALIGN.RIGHT, space_after=3)
    para(tf, "github.com/UbaJaz/Digital_Minds_Research", size=11.5, font=MONO, color=MUTED,
         align=PP_ALIGN.RIGHT)

    footer(s, 1, "Beaten by a Cheap Surface Classifier")


def slide_design(prs):
    s = blank(prs)
    eyebrow(s, "The question, and the design that can answer it")
    heading(s, "Privileged access means beating a cheap observer", size=31)

    tf = textbox(s, MARGIN, Inches(1.58), Inches(11.9), Inches(0.6))
    rich(tf, [
        ("Two confounds block any black-box test. ", 14.5, INK, True, SANS),
        ("The self predictor is usually also the strongest model in the comparison, and a hidden "
         "property may simply be legible in the text. We remove the first by construction and turn "
         "the second into a measurement.", 14.5, MUTED, False, SANS),
    ], first=True, space_after=0)

    # --- the three models
    models = [
        ("M", "TARGET", "Llama-3.1-70B", "generates a column;\npredicts both", GREEN),
        ("N", "SIBLING", "Hermes-3-70B", "same pretrained base;\ndifferent post-training", BLUE),
        ("F", "FAR OBSERVER", "Mistral-Small-24B", "different organisation,\nbase and family", MUTED),
    ]
    cy = Inches(2.48)
    for i, (k, role, name, note, col) in enumerate(models):
        x = MARGIN + i * Inches(2.62)
        rect(s, x, cy, Inches(2.42), Inches(1.66), fill=PANEL, line=RULE)
        tf = textbox(s, x + Inches(0.16), cy + Inches(0.14), Inches(2.1), Inches(0.34))
        rich(tf, [(k + "   ", 20, col, True, MONO), (role, 10, DIM, True, MONO)],
             first=True, space_after=0)
        tf = textbox(s, x + Inches(0.16), cy + Inches(0.62), Inches(2.1), Inches(0.95))
        para(tf, name, size=13.5, color=INK, bold=True, first=True, space_after=4)
        for ln in note.split("\n"):
            para(tf, ln, size=11.5, color=MUTED, space_after=0, line=1.16)

    # --- the crossed grid
    gx = MARGIN + Inches(8.15)
    tf = textbox(s, gx, cy - Inches(0.3), Inches(4.4), Inches(0.28))
    para(tf, "CROSSED 2×2 — 24 CELLS", size=10.5, font=MONO, color=AMBER, first=True,
         bold=True, spacing=1.2)
    cw, ch = Inches(1.28), Inches(0.4)
    tf = textbox(s, gx + Inches(1.05), cy - Inches(0.02), cw * 2 + Inches(0.1), Inches(0.26))
    para(tf, "M's texts      N's texts", size=10.5, font=MONO, color=DIM, first=True)
    for r_i, pred in enumerate(["M", "N", "F"]):
        y = cy + Inches(0.26) + r_i * (ch + Inches(0.06))
        tf = textbox(s, gx, y + Inches(0.08), Inches(1.0), Inches(0.28))
        para(tf, f"{pred} predicts", size=11, font=MONO, color=MUTED, first=True)
        for c_i, colname in enumerate(["M", "N"]):
            is_self = pred == colname
            box = rect(s, gx + Inches(1.05) + c_i * (cw + Inches(0.1)), y, cw, ch,
                       fill=RGBColor(0x1F, 0x33, 0x24) if is_self else PANEL2,
                       line=GREEN if is_self else RULE)
            btf = box.text_frame
            btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            para(btf, f"{pred}→{colname}" + ("  self" if is_self else ""), size=11,
                 font=MONO, color=GREEN if is_self else MUTED, bold=is_self, first=True,
                 align=PP_ALIGN.CENTER)
    # baseline strip, directly under the grid
    dy = cy + Inches(0.26) + 3 * (ch + Inches(0.06)) + Inches(0.06)
    rect(s, gx, dy, Inches(3.71), Inches(0.42), fill=PANEL, line=AMBER)
    dtf = textbox(s, gx + Inches(0.14), dy + Inches(0.1), Inches(3.5), Inches(0.3))
    rich(dtf, [("D  ", 12, AMBER, True, MONO),
               ("surface baseline — the cheap third party", 11.5, INK, False, SANS)],
         first=True, space_after=0)

    tf = textbox(s, MARGIN, Inches(4.62), Inches(11.9), Inches(1.5))
    rich(tf, [
        ("Capability cancels. ", 14, GREEN, True, SANS),
        ("An additive competence edge appears in both of M's cells, so it drops out of the "
         "interaction  ", 14, MUTED, False, SANS),
        ("(M→M − N→M) − (M→N − N→N)", 13, INK, True, MONO),
        (".", 14, MUTED, False, SANS),
    ], first=True, space_after=6)
    rich(tf, [
        ("Leakage becomes the variable. ", 14, AMBER, True, SANS),
        ("Four stimulus constructions on one shared 200-prompt pool, spanning surface baselines "
         "0.54–0.85, with D fit per target column and gated on ", 14, MUTED, False, SANS),
        ("before", 14, INK, True, SANS),
        (" main data is collected.", 14, MUTED, False, SANS),
    ], space_after=6)
    rich(tf, [
        ("The criterion. ", 14, BLUE, True, SANS),
        ("Privileged access requires beating an equal-or-lower-cost observer reading the same "
         "text — not merely beating chance.", 14, MUTED, False, SANS),
    ], space_after=0)

    rect(s, MARGIN, Inches(6.18), BODY_W, Inches(0.52), fill=PANEL2)
    stf = textbox(s, MARGIN + Inches(0.2), Inches(6.31), BODY_W - Inches(0.4), Inches(0.3))
    rich(stf, [
        ("Ground truth is constructed, not elicited. ", 12.5, INK, True, SANS),
        ("The label lives where the prediction code structurally cannot import it — no predictor, "
         "including Self, ever sees it.", 12.5, MUTED, False, SANS),
    ], first=True, space_after=0)

    footer(s, 2, "Question and design")


def slide_leakage(prs):
    s = blank(prs)
    eyebrow(s, "Finding 1  ·  the pilot, five persona pairs, ten generator columns")
    heading(s, "Persona prediction tracks a cheap surface signal")

    picture(s, "fig1_self_vs_surface.png", y=Inches(1.85), height=Inches(4.35),
            cx=Inches(3.85))

    x = Inches(7.55)
    stats = [("r = +0.71", "correlation between self-prediction accuracy and the 21-feature "
                           "surface/textual baseline, across ten columns", AMBER),
             ("6 / 10", "columns where the surface baseline matches or beats the model at "
                        "reading its own persona", AMBER),
             ("0.325", "surface baseline on VO-D's M column once style is equalised — and the "
                       "model falls to 0.500 with it", BLUE)]
    y = Inches(1.86)
    for v, k, col in stats:
        tf = textbox(s, x, y, Inches(5.0), Inches(1.05))
        para(tf, v, size=28, font=MONO, color=col, bold=True, first=True, space_after=2)
        para(tf, k, size=12.5, color=MUTED, line=1.22, space_after=0)
        y += Inches(1.08)

    tf = textbox(s, x, Inches(5.08), Inches(5.05), Inches(0.42))
    para(tf, "21 textual features: 18 structural/style + 3 preregistered persona-linked lexical rates",
         size=11.5, color=DIM, first=True, line=1.2)

    rect(s, x, Inches(5.52), Inches(5.05), Inches(1.22), fill=PANEL, line=RED)
    tf = textbox(s, x + Inches(0.18), Inches(5.64), Inches(4.7), Inches(1.0))
    rich(tf, [("Caveat.  ", 12.5, RED, True, SANS),
              ("VO-D/VO-E also changed the behavioural expression of the hidden property, so they "
               "are not clean causal isolation of style: the two personas largely converged on the "
               "same recommendation.", 12.5, INK, False, SANS)],
         first=True, space_after=0, line=1.2)

    tf = textbox(s, MARGIN, Inches(6.42), Inches(6.6), Inches(0.4))
    para(tf, "Points at or below the diagonal: a surface/textual baseline matches or beats the model.",
         size=12, color=DIM, first=True, line=1.2)

    footer(s, 3, "Surface leakage")


def slide_crossed(prs):
    s = blank(prs)
    eyebrow(s, "Finding 2  ·  24 cells  ·  9,269 scored trials  ·  zero malformed")
    heading(s, "No positive raw M-target advantage — and one positive interaction",
            size=30)

    rows = [
        ("Stimulus set", "Surface D\n(M / N col)", "Self-advantage\nM→M − N→M",
         "Capability-controlled\ninteraction", None),
        ("VO-D  style-equalised", "0.55 / 0.54", "+0.000 [−0.015, +0.015]",
         "−0.006 [−0.033, +0.021]", None),
        ("VO-B  original", "0.65 / 0.75", "+0.000 [−0.033, +0.035]",
         "+0.005 [−0.040, +0.050]", None),
        ("VO-A  original", "0.66 / 0.75", "+0.020 [−0.015, +0.056]",
         "−0.030 [−0.079, +0.018]", None),
        ("VO-C  leakiest", "0.69 / 0.85", "−0.033 [−0.058, −0.008]",
         "+0.089 [+0.048, +0.131]", AMBER),
    ]
    top = Inches(1.78)
    tbl = s.shapes.add_table(len(rows), 4, MARGIN, top, Inches(9.0), Inches(2.5)).table
    for w, col in zip((Inches(2.35), Inches(1.55), Inches(2.55), Inches(2.55)), tbl.columns):
        col.width = w
    for ri, row in enumerate(rows):
        hi = row[4]
        for ci, val in enumerate(row[:4]):
            cell = tbl.cell(ri, ci)
            cell.text = val
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(10.5 if ri == 0 else 12)
                    r.font.name = MONO if (ri and ci) else SANS
                    r.font.bold = (ri == 0) or bool(hi)
                    r.font.color.rgb = DIM if ri == 0 else (hi or INK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (PANEL2 if hi else (PANEL if ri else GROUND))

    # right-hand call-out
    card = rect(s, Inches(10.05), top, Inches(2.55), Inches(2.5), fill=PANEL, line=AMBER)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.16)
    para(tf, "+0.089", size=30, font=MONO, color=AMBER, bold=True, first=True, space_after=2)
    para(tf, "[+0.048, +0.131]", size=11.5, font=MONO, color=MUTED, space_after=7)
    para(tf, "The originally preregistered interaction — positive, on the leakiest set in the "
             "study. We report it rather than bury it.", size=12, color=INK, line=1.22,
         space_after=0)

    tf = textbox(s, MARGIN, Inches(4.52), Inches(11.9), Inches(1.3))
    rich(tf, [
        ("Not one construction shows a positive M-target self-advantage whose interval excludes "
         "zero. ", 15, INK, True, SANS),
        ("The single significant value on that contrast is ", 15, MUTED, False, SANS),
        ("negative", 15, RED, True, SANS),
        (" — and on VO-C's M column the self model is the worst of the three predictors of its "
         "own output (M→M 0.603, below N→M 0.636 and F→M 0.628).", 15, MUTED, False, SANS),
    ], first=True, space_after=7)
    rich(tf, [
        ("Positive interaction, but predictor-by-column differences prevent it from uniquely "
         "identifying privileged access. ", 14, AMBER, True, SANS),
        ("It is positive because M ", 14, MUTED, False, SANS),
        ("under", 14, INK, True, SANS),
        ("-performs on N's column, not because it over-performs on its own — and the estimator "
         "cancels only the additive part.", 14, MUTED, False, SANS),
    ], space_after=0)

    strip = rect(s, MARGIN, Inches(6.15), BODY_W, Inches(0.62), fill=PANEL2)
    stf = strip.text_frame
    stf.margin_left = Inches(0.2)
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    rich(stf, [
        ("Four stimulus constructions  ·  shared 200-prompt pool", 13, INK, True, SANS),
        ("          Preregistered floor 500/cell; achieved ~400; VO-D N = 323. Prompt-clustered "
         "bootstrap throughout.", 11, DIM, False, SANS),
    ], first=True, space_after=0)

    footer(s, 4, "Crossed design")


def slide_selfpred(prs):
    s = blank(prs)
    eyebrow(s, "Finding 3  ·  “which of these two replies would you produce?”  ·  391 pairs")
    heading(s, "Hermes-3 can predict itself. A cheaper observer still does better.", size=30)

    # three comparison columns — length is the matched observer; 21-feature is a cost criterion
    bars = [
        ("0.719", "Hermes-3 self-prediction", "MODEL RESULT",
         "balanced accuracy [0.675, 0.762]\ndiscrimination +0.437", GREEN, 0.719),
        ("0.808", "Length-only observer", "MATCHED COMPARATOR",
         "“pick the longer reply”\none feature, no training", BLUE, 0.808),
        ("0.831", "21-feature surface classifier", "COST CRITERION",
         "supervised single-text\nauthorship labelling", AMBER, 0.831),
    ]
    lo, hi = 0.50, 0.90
    bw = Inches(3.78)
    track_w = Inches(3.38)
    cy = Inches(1.72)
    for i, (v, title, role, note, col, val) in enumerate(bars):
        x = MARGIN + i * (bw + Inches(0.22))
        # matched length gets a coloured border; criterion stays quieter
        border = BLUE if i == 1 else (GREEN if i == 0 else RULE)
        rect(s, x, cy, bw, Inches(2.28), fill=PANEL, line=border)
        tf = textbox(s, x + Inches(0.2), cy + Inches(0.1), bw - Inches(0.4), Inches(1.6))
        para(tf, role, size=9.5, font=MONO, color=col, bold=True, first=True, space_after=1,
             spacing=1.0)
        para(tf, title, size=13, color=INK, bold=True, space_after=2)
        para(tf, v, size=36, font=MONO, color=col, bold=True, space_after=2)
        for ln in note.split("\n"):
            para(tf, ln, size=11.5, color=MUTED, space_after=0, line=1.16)
        # bar on an explicitly labelled 0.50-0.90 track
        ty = cy + Inches(1.72)
        rect(s, x + Inches(0.2), ty, track_w, Inches(0.15), fill=PANEL2)
        rect(s, x + Inches(0.2), ty, Emu(int(track_w * (val - lo) / (hi - lo))), Inches(0.15),
             fill=col)
        tf2 = textbox(s, x + Inches(0.2), ty + Inches(0.19), track_w, Inches(0.2))
        para(tf2, "0.50 chance", size=9.5, font=MONO, color=DIM, first=True, space_after=0)
        tf3 = textbox(s, x + Inches(0.2), ty + Inches(0.19), track_w, Inches(0.2))
        para(tf3, "0.90", size=9.5, font=MONO, color=DIM, first=True, space_after=0,
             align=PP_ALIGN.RIGHT)

    tf = textbox(s, MARGIN, Inches(4.06), Inches(11.9), Inches(0.5))
    rich(tf, [
        ("The 0.831 classifier is a different evaluation procedure; the only matched comparator "
         "is the length rule ", 12.5, MUTED, False, SANS),
        ("(paired difference +0.095 [+0.036, +0.155], McNemar p = 0.0018)", 12.5, INK, False, MONO),
        (".", 12.5, MUTED, False, SANS),
    ], first=True, space_after=0)

    # the residual — visually distinct
    rect(s, MARGIN, Inches(4.6), Inches(7.7), Inches(1.32), fill=RGBColor(0x1F, 0x33, 0x24),
         line=GREEN)
    tf = textbox(s, MARGIN + Inches(0.2), Inches(4.73), Inches(7.3), Inches(1.1))
    para(tf, "Hermes still discriminates where the length cue points away from its own reply:",
         size=13, color=INK, bold=True, first=True, space_after=2)
    para(tf, "+0.381 [+0.188, +0.566]", size=17, font=MONO, color=GREEN, bold=True, space_after=3)
    para(tf, "75 pairs where a pure length strategy is actively wrong. This study cannot name "
             "the residual.", size=11.5, color=MUTED, line=1.16, space_after=0)

    rect(s, Inches(8.62), Inches(4.6), Inches(4.0), Inches(1.32), fill=PANEL2, line=RULE)
    tf = textbox(s, Inches(8.8), Inches(4.73), Inches(3.64), Inches(1.1))
    para(tf, "Llama-3.1: no self-prediction", size=13, color=INK, bold=True, first=True,
         space_after=3)
    para(tf, "Discrimination −0.107 [−0.166, −0.048]; 89.7% “A”. One model's ability, not a "
             "property of language models.", size=12, color=MUTED, line=1.18, space_after=0)

    rect(s, MARGIN, Inches(6.06), BODY_W, Inches(0.48), fill=GROUND, line=AMBER)
    tf = textbox(s, MARGIN + Inches(0.22), Inches(6.17), BODY_W - Inches(0.44), Inches(0.3))
    rich(tf, [("Self-prediction is possible. ", 15.5, AMBER, True, SANS),
              ("Privileged self-access is not thereby demonstrated.", 15.5, INK, True, SANS)],
         first=True, space_after=0)

    tf = textbox(s, MARGIN, Inches(6.62), BODY_W, Inches(0.28))
    para(tf, "post hoc analyses  ·  0.719 and 0.831 are different evaluation procedures, and no "
             "statistical test is run between them", size=10.5, font=MONO, color=DIM, first=True)

    footer(s, 5, "Self-prediction")


def slide_future(prs):
    s = blank(prs)
    eyebrow(s, "Where this goes next")
    heading(s, "A decision tree, not a schedule")

    tf = textbox(s, MARGIN, Inches(1.68), Inches(11.9), Inches(0.5))
    rich(tf, [
        ("One question is left open by our own data: ", 15.5, MUTED, False, SANS),
        ("what is the model-specific residual that survives control of cheap surface cues, and can "
         "a behavioural test be built in which a positive privileged-access result would be "
         "identifiable?", 15.5, INK, True, SANS),
    ], first=True, space_after=0)

    stages = [
        ("1", "DISSOCIATE", GREEN,
         "Self-preference vs self-prediction",
         "Re-run the probe on the same pairs under two questions — “which would you produce?” "
         "against “which is better?”. A residual as large under the quality question reads as "
         "self-preference; one specific to the prediction framing is the discriminating outcome. "
         "Needs a behavioural manipulation check run before collection."),
        ("2", "AUDIT", BLUE,
         "Apply the controls to existing claims",
         "Take the leakage gate and the response-bias check to published behavioural "
         "introspection results, and ask how many survive controls this cheap. No new model runs; "
         "it tests whether the framework generalises beyond our stimuli."),
        ("3", "TEST", AMBER,
         "Stronger causal ground truth — conditional",
         "Only if a residual survives Stage 1 and audited effects do not dissolve: a training-"
         "relationship ladder against our one-lineage limit, activation steering or an "
         "independently planted property so ground truth is verified rather than assumed, and an "
         "incremental-validity test against an observer's features."),
    ]
    cw = Inches(3.86)
    cy = Inches(2.42)
    for i, (num, kicker, col, title, body) in enumerate(stages):
        x = MARGIN + i * (cw + Inches(0.22))
        rect(s, x, cy, cw, Inches(3.05), fill=PANEL, line=RULE)
        rect(s, x, cy, cw, Inches(0.06), fill=col)
        # fixed rows, so the three cards align exactly rather than flowing independently
        tf = textbox(s, x + Inches(0.2), cy + Inches(0.24), cw - Inches(0.4), Inches(0.34))
        rich(tf, [(num + "   ", 20, col, True, MONO), (kicker, 12.5, col, True, MONO)],
             first=True, space_after=0)
        tf = textbox(s, x + Inches(0.2), cy + Inches(0.72), cw - Inches(0.4), Inches(0.5))
        para(tf, title, size=14, color=INK, bold=True, first=True, space_after=0, line=1.15)
        tf = textbox(s, x + Inches(0.2), cy + Inches(1.3), cw - Inches(0.4), Inches(1.6))
        para(tf, body, size=11.5, color=MUTED, line=1.28, first=True, space_after=0)
        if i < 2:
            tfa = textbox(s, x + cw + Inches(0.02), Inches(3.75), Inches(0.2), Inches(0.4))
            para(tfa, "→", size=17, color=DIM, first=True, align=PP_ALIGN.CENTER)

    tf = textbox(s, MARGIN, Inches(5.68), Inches(11.9), Inches(0.5))
    para(tf, "If Stage 1 dissolves the residual, that is the result and Stage 3 does not run. The "
             "bottleneck is breadth and experimental design, not implementation.",
         size=12.5, color=DIM, first=True, line=1.24)

    rect(s, MARGIN, Inches(6.12), BODY_W, Inches(0.58), fill=PANEL2, line=AMBER)
    tf = textbox(s, MARGIN + Inches(0.24), Inches(6.26), BODY_W - Inches(0.48), Inches(0.34))
    para(tf, "Can behavioural self-prediction ever provide evidence unavailable to an external "
             "observer?", size=16.5, color=AMBER, bold=True, first=True, space_after=0)

    footer(s, 6, "Future work")


def build() -> None:
    prs = new_deck()
    slide_title(prs)
    slide_design(prs)
    slide_leakage(prs)
    slide_crossed(prs)
    slide_selfpred(prs)
    slide_future(prs)
    prs.save(OUT)
    print(f"-> {OUT.name}  ({OUT.stat().st_size / 1024:.0f} KB, {len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
